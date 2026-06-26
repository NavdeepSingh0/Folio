import fitz  # PyMuPDF
from pptx import Presentation
import io

def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    """Extracts text from a given file (PDF or PPTX)."""
    text = ""
    
    if filename.lower().endswith(".pdf"):
        # Extract from PDF
        pdf_document = fitz.open(stream=file_bytes, filetype="pdf")
        for page_num in range(pdf_document.page_count):
            page = pdf_document.load_page(page_num)
            text += page.get_text() + "\n"
        pdf_document.close()
        
    elif filename.lower().endswith(".pptx"):
        # Extract from PPTX
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
            
    else:
        raise ValueError("Unsupported file format. Please upload PDF or PPTX.")
        
    return text.strip()
