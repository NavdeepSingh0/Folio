import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import re

def extract_structured_text(file_bytes: bytes, filename: str) -> tuple[str, int]:
    """Extracts structured text from a PDF or PPTX and returns (text, page_count)."""
    if filename.lower().endswith(".pdf"):
        return _extract_pdf_structure(file_bytes)
    elif filename.lower().endswith(".pptx"):
        return _extract_pptx_structure(file_bytes)
    else:
        raise ValueError("Unsupported file format. Please upload PDF or PPTX.")

def _extract_pptx_structure(file_bytes: bytes) -> tuple[str, int]:
    prs = Presentation(io.BytesIO(file_bytes))
    markdown_content = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        is_quiz_slide = False
        
        # Determine if this is a quiz slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.lower()
                if "quiz" in text or "discussion" in text or "questions" in text or "think" in text:
                    is_quiz_slide = True
                    break
        
        if is_quiz_slide:
            slide_text.append("## Self Test")

        for shape in slide.shapes:
            if getattr(shape, "is_placeholder", False):
                # Try to detect if it's a title (type 1 is usually TITLE)
                # But it's safer to just check if it's the first text placeholder
                if shape.has_text_frame and shape.text.strip():
                    if hasattr(shape, "placeholder_format") and shape.placeholder_format.type in (1, 3): # TITLE or SUBTITLE
                        if not is_quiz_slide:
                            slide_text.append(f"# {shape.text.strip()}")
                        continue

            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    # Check for Important boxes
                    lower_text = text.lower()
                    if lower_text.startswith(("definition", "important", "warning", "note", "tip", "example")):
                        slide_text.append(f"> {text}")
                    else:
                        # Process bullets vs normal text
                        for paragraph in shape.text_frame.paragraphs:
                            p_text = paragraph.text.strip()
                            if p_text:
                                if paragraph.level > 0 or (hasattr(paragraph, 'level') and paragraph.level == 0 and len(shape.text_frame.paragraphs) > 1):
                                    # It's likely a list
                                    indent = "  " * getattr(paragraph, 'level', 0)
                                    slide_text.append(f"{indent}- {p_text}")
                                else:
                                    slide_text.append(p_text)

            elif shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    row_data = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    slide_text.append("| " + " | ".join(row_data) + " |")
                    if row_idx == 0:
                        slide_text.append("|" + "|".join(["---"] * len(row.cells)) + "|")
                slide_text.append("\n")

            elif shape.shape_type == getattr(MSO_SHAPE_TYPE, 'PICTURE', 13):
                slide_text.append(f"[Figure: Image on Slide {i + 1}]")

        markdown_content.append("\n".join(slide_text))
        markdown_content.append("\n---\n")

    return "\n".join(markdown_content), len(prs.slides)

def _extract_pdf_structure(file_bytes: bytes) -> tuple[str, int]:
    pdf_document = fitz.open(stream=file_bytes, filetype="pdf")
    markdown_content = []

    for page_num in range(pdf_document.page_count):
        page = pdf_document.load_page(page_num)
        page_text = []

        # 1. Extract Images
        image_list = page.get_images(full=True)
        if image_list:
            for img_index, img in enumerate(image_list):
                page_text.append(f"[Figure: Image {img_index + 1} on Page {page_num + 1}]")

        # 2. Extract Tables
        try:
            tables = page.find_tables()
            for table in tables:
                rows = table.extract()
                if rows and len(rows) > 0:
                    for r_idx, row in enumerate(rows):
                        # Clean cells
                        cleaned_row = [str(cell).strip().replace('\n', ' ') if cell is not None else "" for cell in row]
                        page_text.append("| " + " | ".join(cleaned_row) + " |")
                        if r_idx == 0:
                            page_text.append("|" + "|".join(["---"] * len(row)) + "|")
                    page_text.append("\n")
        except Exception as e:
            print(f"Failed to parse table on page {page_num}: {e}")

        # 3. Extract Text with rough heading detection
        blocks = page.get_text("dict")["blocks"]
        for b in blocks:
            if b['type'] == 0:  # text block
                for line in b["lines"]:
                    for span in line["spans"]:
                        text = span["text"].strip()
                        if not text:
                            continue
                        
                        size = span["size"]
                        # Rough heuristic: size > 14 is a heading
                        if size > 14:
                            page_text.append(f"## {text}")
                        else:
                            # Check for Important boxes
                            lower_text = text.lower()
                            if lower_text.startswith(("definition", "important", "warning", "note", "tip", "example")):
                                page_text.append(f"> {text}")
                            else:
                                page_text.append(text)
        
        markdown_content.append("\n\n".join(page_text))
        markdown_content.append("\n---\n")

    page_count = pdf_document.page_count
    pdf_document.close()
    
    # 4. Reference Grouping
    full_text = "\n".join(markdown_content)
    
    # Simple regex to group references if they are scattered, but realistically
    # they are usually at the end. We will let the LLM handle semantic grouping,
    # but the structural extraction at least preserves them cleanly.
    
    return full_text, page_count
